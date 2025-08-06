from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Create a new presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

# Sample data for the donut chart
chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Sales', (20.4, 30.6, 15.2, 33.8))

# Add chart to slide
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
).chart

# Control the donut width by adjusting inner_radius
# Values range from 0.0 to 1.0
# 0.0 = no hole (pie chart)
# 1.0 = maximum hole (very thin ring)

# Thin donut (wide ring)
chart.plots[0].inner_radius = 0.2  # 20% inner radius = thick donut

# Medium donut
# chart.plots[0].inner_radius = 0.5  # 50% inner radius = medium donut

# Thick donut (thin ring)  
# chart.plots[0].inner_radius = 0.8  # 80% inner radius = thin donut

# Optional: Customize other chart properties
chart.has_legend = True
# chart.legend.position = XL_LEGEND_POSITION.BOTTOM

# Save the presentation
prs.save('donut_chart_example.pptx')

print("Donut chart created with custom width!")

# Additional examples of different widths:
def create_donut_with_width(inner_radius, filename):
    """Create a donut chart with specified inner radius (width)"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    chart_data = CategoryChartData()
    chart_data.categories = ['Category A', 'Category B', 'Category C', 'Category D']
    chart_data.add_series('Values', (25, 35, 20, 20))
    
    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart
    
    # Set the inner radius (controls donut width)
    chart.plots[0].inner_radius = inner_radius
    
    # Add title
    title = slide.shapes.title
    title.text = f"Donut Chart (Inner Radius: {inner_radius})"
    
    prs.save(filename)
    print(f"Created {filename} with inner radius: {inner_radius}")

# Create examples with different widths
create_donut_with_width(0.1, 'thick_donut.pptx')    # Very thick donut
# create_donut_with_width(0.3, 'medium_donut.pptx')   # Medium donut  
# create_donut_with_width(0.6, 'thin_donut.pptx')     # Thin donut
# create_donut_with_width(0.9, 'very_thin_donut.pptx') # Very thin donut