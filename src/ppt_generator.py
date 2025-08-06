import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_AXIS_CROSSES, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime


class PptGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.current_slide = None
    

    def create_blank_slide(self, slide_title: str | None = None):
        """Create a blank slide"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        if slide_title:
            slide.shapes.title.text = slide_title
            slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            # slide.shapes.title.width = Inches(7)
        self.current_slide = slide

    def add_image_header_footer_to_all_slides(self, image_path: str):
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
    
        img_width=Inches(1.85)
        img_height=Inches(0.17)
        left_position = slide_width - img_width - Inches(0.3)
        top_position = Inches(0.3)
        
        for slide in self.prs.slides:
            # Header
            slide.shapes.add_picture(
                image_path,
                left_position,
                top_position,
                width=img_width,
                height=img_height
            )
            
            # Footer
            footer = slide.shapes.add_textbox(
                left=(slide_width - Inches(6)) / 2,
                top=slide_height - Inches(0.5),
                width=Inches(6),
                height=Inches(0.3)
            )
            footer_frame = footer.text_frame
            footer_frame.text = f"優才資源中心有限公司{datetime.now().year}年DSE考生問卷調查 考生未來勞動力供應剖析"
            footer_para = footer_frame.paragraphs[0]
            footer_para.font.size = Pt(12)
            footer_para.font.color.rgb = RGBColor(100, 100, 100)
            footer_para.alignment = PP_ALIGN.CENTER
    
    def add_textbox(
        self,
        text: list[str] | str,
        x: float = 1,
        y: float = 2,
        cx: float = 8,
        cy: float = 5,
        font_size: int = 18,
    ):
        """Add a textbox to the current slide"""
        if isinstance(text, list):
            text = "\n".join(text)

        if self.current_slide is None:
            self.create_blank_slide()

        textbox = self.current_slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cx), Inches(cy))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        
        p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
        p.alignment = PP_ALIGN.LEFT                     


    def add_bar_chart(
        self,
        data: pd.DataFrame,
        category_column: str,
        value_columns: list[str],
        title: str = None,
        has_legend: bool = True,
        legend_position: int = 2,  
        to_percentage: bool = False,
        hide_y_axis: bool = False,
        opposite_tick_labels: bool = False,
        reserve_value_axis: bool = False,
        font_size: int = 14,
        small_title: bool = False,
        x: float = 1,
        y: float = 2,
        cx: float = 8,
        cy: float = 5,
        horizontal: bool = False,
    ):

        if category_column not in data.columns :
            print(f"Columns {category_column} not found in data")
            return

        for value_column in value_columns:
            if value_column not in data.columns:
                print(f"Value column {value_column} not found in data")
                return

        if self.current_slide is None:
            self.create_blank_slide()

        # Prepare chart data
        chart_data = CategoryChartData()
        chart_data.categories = data[category_column].tolist()

        for value_column in value_columns:
            chart_data.add_series(value_column, data[value_column].tolist())

        # Choose chart type based on orientation
        chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
        
            
        # Add chart to slide
        x, y, cx, cy = Inches(x), Inches(y), Inches(cx), Inches(cy)
        chart = self.current_slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart

        # Set chart title
        if title is None:
            chart.has_title = False
        else:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            if small_title:
                title_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
                title_run.font.size = Pt(12)
                title_run.font.bold = False
                title_run.font.name = 'Calibri'

        chart.has_legend = has_legend
        if has_legend:
            chart.legend.position = legend_position
            chart.legend.font.size = Pt(font_size)

            if len(value_columns) == 1 or legend_position == XL_LEGEND_POSITION.BOTTOM:
                chart.legend.include_in_layout = False
            else:
                chart.legend.include_in_layout = True

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.show_value = True
        data_labels.show_category_name = False


        if to_percentage:
            data_labels.number_format = '0.0%'  # Format as percentage with two decimal places
            chart.value_axis.tick_labels.number_format = '0%'

        if hide_y_axis:
            chart.value_axis.visible = False
        
        # Move x-axis labels to top
        if opposite_tick_labels:
            chart.category_axis.crosses = XL_AXIS_CROSSES.MAXIMUM
            chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH
            chart.category_axis.reverse_order = True

        if reserve_value_axis:
            chart.value_axis.reverse_order = True

        # Change axis font size
        chart.category_axis.tick_labels.font.size = Pt(font_size)
        chart.value_axis.tick_labels.font.size = Pt(font_size)
        data_labels.font.size = Pt(font_size)

        chart.value_axis.has_major_gridlines = False
        chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart.category_axis.minor_tick_mark = XL_TICK_MARK.NONE
        chart.value_axis.minor_tick_mark = XL_TICK_MARK.NONE


    def add_pie_chart(
        self,
        data: pd.DataFrame,
        category_column: str,
        value_column: str,
        to_percent: bool = False,
        title: str | None = None,
        has_legend: bool = True,
        legend_position: int = 2,  # 2 for bottom
        sort: bool = True,
        font_size: int = 14,
        max_categories: int = 8,
        x: float = 1,
        y: float = 2,
        cx: float = 8,
        cy: float = 5,
    ):

        if category_column not in data.columns or value_column not in data.columns:
            print(f"Columns {category_column} or {value_column} not found in data")
            return

        if self.current_slide is None:
            self.create_blank_slide(title)

        pie_data = data.groupby(category_column)[value_column].sum().reset_index()
        if sort:
            pie_data = pie_data.sort_values(value_column, ascending=False)

        # Group smaller categories into 'Other'
        if len(pie_data) > max_categories:
            top = pie_data.iloc[:max_categories]
            other = pd.DataFrame({
                category_column: ['Other'],
                value_column: [pie_data.iloc[max_categories:][value_column].sum()]
            })
            pie_data = pd.concat([top, other], ignore_index=True)

        chart_data = CategoryChartData()
        chart_data.categories = pie_data[category_column].tolist()
        chart_data.add_series(value_column, pie_data[value_column].tolist())

        x, y, cx, cy = Inches(x), Inches(y), Inches(cx), Inches(cy)
        chart = self.current_slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart
        
        # Set chart title
        if title is None:
            chart.has_title = False
        else:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        chart.has_legend = has_legend
        if has_legend:
            chart.legend.position = legend_position
            chart.legend.include_in_layout = False 
            chart.legend.font.size = Pt(font_size)

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.show_value = True
        data_labels.font.size = Pt(font_size)

        if to_percent:
            data_labels.show_percentage = True
            data_labels.show_value = False



    def add_donut_chart(
        self,
        data: pd.DataFrame,
        category_column: str,
        value_column: str,
        to_percent: bool = False,
        sort: bool = True,
        title: str = "Donut Chart",
        has_legend: bool = True,
        legend_position: int = 2,  # 2 for bottom 3 for right
        has_data_labels: bool = False,
        data_labels_outside: bool = False,  # NEW: Position labels outside chart sections
        font_size: int = 14,
        max_categories: int = 8,
        small_title: bool = False,
        x: float = 1,
        y: float = 2,
        cx: float = 8,
        cy: float = 5,
    ):
        """
        Add a donut chart to the current slide.
        
        Parameters:
        -----------
        data_labels_outside : bool, default False
            If True, positions data labels outside the chart sections.
            If False, positions data labels inside the sections (default behavior).
        """
        if category_column not in data.columns or value_column not in data.columns:
            print(f"Columns {category_column} or {value_column} not found in data")
            return

        if self.current_slide is None:
            self.create_blank_slide(title)

        pie_data = data.groupby(category_column)[value_column].sum().reset_index()
        if sort:
            pie_data = pie_data.sort_values(value_column, ascending=False)

        # Group smaller categories into 'Other'
        if len(pie_data) > max_categories:
            top = pie_data.iloc[:max_categories]
            other = pd.DataFrame({
                category_column: ['Other'],
                value_column: [pie_data.iloc[max_categories:][value_column].sum()]
            })
            pie_data = pd.concat([top, other], ignore_index=True)

        chart_data = CategoryChartData()
        chart_data.categories = pie_data[category_column].tolist()

        chart_data.add_series(value_column, pie_data[value_column].tolist())

        x, y, cx, cy = Inches(x), Inches(y), Inches(cx), Inches(cy)
        chart = self.current_slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
        ).chart

        # Set chart title
        chart.has_title = True
        chart.chart_title.text_frame.text = title if title else category_column
        if small_title:
            title_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
            title_run.font.size = Pt(12)
            title_run.font.bold = False
            title_run.font.name = 'Calibri'

        chart.has_legend = has_legend
        if has_legend:
            chart.legend.position = legend_position
            chart.legend.include_in_layout = True
            chart.legend.font.size = Pt(font_size)

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.show_value = True
        data_labels.font.size = Pt(font_size)

        if has_data_labels:
            data_labels.show_category_name = True
        
        if to_percent:
            data_labels.show_percentage = True
            data_labels.show_value = False

        chart.doughnut_hole_size = 90


    def add_table(
        self,
        data: pd.DataFrame,
        rows: int | None = None,
        cols: int | None = None,
        index: bool = True,
        x: float = 1,
        y: float = 2,
        cx: float = 8,
        cy: float = 5,
        font_size: int = 12,
    ):
        """Add a table to the current slide"""
        if self.current_slide is None:
            self.create_blank_slide()
        
        if rows is None or cols is None:
            rows = len(data) + 1  # +1 for header row
            cols = len(data.columns) + 1 if index else len(data.columns) 



        table = self.current_slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(cx), Inches(cy)).table


        # Fill the table with data, first column as index
        for i, col in enumerate(data.columns):
            cell = table.cell(0, i + 1 if index else i)
            cell.text = col

        if index:
            for j in range(1, rows):
                if j - 1 >= len(data):
                    break
                index_cell = table.cell(j, 0)
                index_cell.text = str(data.index[j - 1])


        for i, col in enumerate(data.columns):
            if index and i + 1 >= cols:
                break
            if not index and i >= cols:
                break
            for j in range(len(data)):
                if j + 1 >= rows:
                    break 
                data_cell = table.cell(j + 1, i + 1 if index else i)
                data_cell.text = str(data.iloc[j][col])

        # Set font size for all cells in a separate loop
        for row in range(rows):
            for col in range(cols):
                cell = table.cell(row, col)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(font_size)

    

    def add_stacked_bar(
        self,
        data: pd.DataFrame,
        category_column: str,
        value_columns: list[str],
        title: str,
        legend_position: int = 2,  
        font_size: int = 12,
        x: float = 1,
        y: float = 2,
        cx: float = 8,
        cy: float = 5,
    ):
        if self.current_slide is None:
            self.create_blank_slide()

        chart_data = CategoryChartData()
        chart_data.categories = data[category_column].tolist()

        for value_column in value_columns:
            chart_data.add_series(value_column, data[value_column].tolist())

        x, y, cx, cy = Inches(x), Inches(y), Inches(cx), Inches(cy)
        chart = self.current_slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data
        ).chart

        # Set chart title
        chart.has_title = True
        chart.chart_title.text_frame.text = title

        chart.has_legend = True
        chart.legend.position = legend_position
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(font_size)
        
        chart.value_axis.tick_labels.number_format = '0%'
        chart.value_axis.maximum_scale = 1
        chart.value_axis.tick_labels.font.size = Pt(font_size)

        chart.category_axis.tick_labels.font.size = Pt(font_size)

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.show_value = True
        data_labels.number_format = '0.0%'
        data_labels.font.size = Pt(font_size)


    def add_img(
        self, 
        path: str, 
        x: float, 
        y: float, 
        cx: float | None = None, 
        cy: float | None = None
    ):
        if self.current_slide is None:
            self.create_blank_slide()

        img_path = path
        if os.path.exists(img_path):
            self.current_slide.shapes.add_picture(
                img_path, Inches(x), Inches(y), 
                width=Inches(cx) if cx else None, height=Inches(cy) if cy else None
            )
        else:
            print(f"Image {img_path} not found.")

    def save(self, path: str):
        try:
            self.prs.save(path)
            print(f"Presentation saved as {path}")
            return True
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return False

def main():
    pass

if __name__ == "__main__":
    main()